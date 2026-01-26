import os
import logging
import uuid
import tempfile
from typing import List, Dict, Any, Tuple
from datetime import datetime
import json
import PyPDF2
from PIL import Image
import cv2
import numpy as np
import pdf2image
import pytesseract
from app.config import settings
from app.services.ai_service import AIService
from app.models import TextBatch

logger = logging.getLogger(__name__)

class FileProcessor:
    def __init__(self):
        self.ai_service = AIService()
        logger.info("FileProcessor initialized")
    
    def _log_operation(self, operation: str, details: Dict[str, Any], is_request: bool = True):
        """Log file processing operations"""
        log_entry = {
            "timestamp": datetime.utcnow().isoformat(),
            "operation": operation,
            "type": "REQUEST" if is_request else "RESPONSE",
            **details
        }
        
        log_type = "REQUEST" if is_request else "RESPONSE"
        logger.info(f"=== FILE PROCESSOR {log_type}: {operation} ===")
        logger.info(f"Details: {json.dumps(log_entry, indent=2, default=str)}")
        logger.info(f"=== END FILE PROCESSOR {log_type} ===")

    def create_batches(self, file_path: str, session_id: str) -> List[TextBatch]:
        """
        Create batches based on page count
        - ≤5 pages: single batch
        - >5 pages: batches of 2-3 pages each
        """
        self._log_operation("CREATE_BATCHES", {
            "file_path": file_path,
            "session_id": session_id
        })
        
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # Determine page count
            if file_ext == '.pdf':
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    page_count = len(pdf_reader.pages)
            else:
                # For images and other files, count as 1 page
                page_count = 1
            
            logger.info(f"Document has {page_count} pages")
            
            batches = []
            
            if page_count <= 5:
                # Single batch for small documents
                batch = TextBatch(
                    batch_id=str(uuid.uuid4()),
                    session_id=session_id,
                    page_range=(1, page_count),
                    text_content="",
                    batch_number=1,
                    total_batches=1
                )
                batches.append(batch)
                logger.info(f"Created 1 batch for {page_count} pages")
            else:
                # Multiple batches for larger documents (2-3 pages each)
                batch_size = 3
                total_batches = (page_count + batch_size - 1) // batch_size
                
                for i in range(total_batches):
                    start_page = i * batch_size + 1
                    end_page = min((i + 1) * batch_size, page_count)
                    
                    batch = TextBatch(
                        batch_id=str(uuid.uuid4()),
                        session_id=session_id,
                        page_range=(start_page, end_page),
                        text_content="",
                        batch_number=i + 1,
                        total_batches=total_batches
                    )
                    batches.append(batch)
                
                logger.info(f"Created {total_batches} batches for {page_count} pages")
            
            self._log_operation("CREATE_BATCHES", {
                "page_count": page_count,
                "batches_created": len(batches),
                "batch_details": [{"batch_id": b.batch_id, "page_range": b.page_range} for b in batches]
            }, is_request=False)
            
            return batches
            
        except Exception as e:
            logger.error(f"Failed to create batches: {str(e)}")
            self._log_operation("CREATE_BATCHES", {
                "error": str(e),
                "success": False
            }, is_request=False)
            # Return single empty batch on error
            return [TextBatch(
                batch_id=str(uuid.uuid4()),
                session_id=session_id,
                page_range=(1, 1),
                text_content="",
                batch_number=1,
                total_batches=1
            )]

    async def extract_text_default(self, file_path: str) -> str:
        """Extract text using default method (direct extraction)"""
        self._log_operation("EXTRACT_TEXT_DEFAULT", {"file_path": file_path})
        
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.pdf':
                text = await self._extract_pdf_text(file_path)
            elif file_ext in ['.jpg', '.jpeg', '.png']:
                text = await self._extract_image_ocr(file_path)
            elif file_ext == '.pptx':
                text = await self._extract_pptx_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            self._log_operation("EXTRACT_TEXT_DEFAULT", {
                "file_path": file_path,
                "text_length": len(text) if text else 0,
                "text_preview": (text[:500] + "...") if text and len(text) > 500 else text,
                "success": True
            }, is_request=False)
            
            return text
                
        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {str(e)}")
            self._log_operation("EXTRACT_TEXT_DEFAULT", {
                "error": str(e),
                "success": False
            }, is_request=False)
            return ""

    async def extract_text_ocr(self, file_path: str, session_id: str = None) -> str:
        """Extract text using OCR method with page-based processing"""
        self._log_operation("EXTRACT_TEXT_OCR", {
            "file_path": file_path,
            "session_id": session_id
        })
        
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            logger.info(f"OCR extraction for file type: {file_ext}")
            
            if file_ext == '.pdf':
                text = await self._extract_pdf_with_ocr_processing(file_path, session_id)
            elif file_ext in ['.jpg', '.jpeg', '.png']:
                text = await self._extract_image_ocr(file_path)
            elif file_ext == '.pptx':
                text = await self._extract_pptx_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            self._log_operation("EXTRACT_TEXT_OCR", {
                "file_path": file_path,
                "text_length": len(text) if text else 0,
                "text_preview": (text[:500] + "...") if text and len(text) > 500 else text,
                "success": True
            }, is_request=False)
            
            return text
                
        except Exception as e:
            logger.error(f"Failed to extract text with OCR from {file_path}: {str(e)}")
            self._log_operation("EXTRACT_TEXT_OCR", {
                "error": str(e),
                "success": False
            }, is_request=False)
            return ""

    async def extract_text_ocr_batched(self, file_path: str, session_id: str) -> List[TextBatch]:
        """
        Extract text using OCR with batching support
        Returns list of TextBatch objects with extracted text
        """
        self._log_operation("EXTRACT_TEXT_OCR_BATCHED", {
            "file_path": file_path,
            "session_id": session_id
        })
        
        try:
            # Create batches first
            batches = self.create_batches(file_path, session_id)
            logger.info(f"Created {len(batches)} batches for OCR extraction")
            
            file_ext = os.path.splitext(file_path)[1].lower()
            logger.info(f"Processing file with extension: {file_ext}")
            
            if file_ext == '.pdf':
                # Extract text for each batch
                for batch in batches:
                    start_page, end_page = batch.page_range
                    logger.info(f"Extracting OCR text for batch {batch.batch_number}: pages {start_page}-{end_page}")
                    
                    batch_text = await self._extract_pdf_pages_ocr(file_path, start_page, end_page)
                    batch.text_content = batch_text
                    
                    logger.info(f"Batch {batch.batch_number} extracted {len(batch_text)} characters")
                    
            elif file_ext in ['.jpg', '.jpeg', '.png']:
                # Single image - single batch
                logger.info("Processing single image with OCR")
                text = await self._extract_image_ocr(file_path)
                batches[0].text_content = text
                logger.info(f"Image OCR extracted {len(text)} characters")
                
            elif file_ext in ['.pptx', '.ppt']:
                # Extract text from presentation
                logger.info("Processing PowerPoint file")
                text = await self._extract_pptx_text(file_path)
                batches[0].text_content = text
                logger.info(f"PPTX extracted {len(text)} characters")
            
            # Log results
            total_text = sum(len(b.text_content) for b in batches)
            self._log_operation("EXTRACT_TEXT_OCR_BATCHED", {
                "batches_processed": len(batches),
                "total_text_length": total_text,
                "batch_details": [
                    {
                        "batch_number": b.batch_number,
                        "page_range": b.page_range,
                        "text_length": len(b.text_content),
                        "text_preview": (b.text_content[:200] + "...") if len(b.text_content) > 200 else b.text_content
                    }
                    for b in batches
                ],
                "success": True
            }, is_request=False)
            
            return batches
            
        except Exception as e:
            logger.error(f"Failed to extract text with OCR batching: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            
            self._log_operation("EXTRACT_TEXT_OCR_BATCHED", {
                "error": str(e),
                "traceback": traceback.format_exc(),
                "success": False
            }, is_request=False)
            
            return []

    async def _extract_pdf_pages_ocr(self, file_path: str, start_page: int, end_page: int) -> str:
        """Extract text from specific PDF pages using OCR"""
        self._log_operation("EXTRACT_PDF_PAGES_OCR", {
            "file_path": file_path,
            "start_page": start_page,
            "end_page": end_page
        })
        
        try:
            logger.info(f"Converting PDF pages {start_page}-{end_page} to images...")
            
            # Convert specific pages to images
            images = pdf2image.convert_from_path(
                file_path,
                first_page=start_page,
                last_page=end_page,
                dpi=200  # Higher DPI for better OCR
            )
            
            logger.info(f"Converted {len(images)} pages to images")
            
            batch_text = ""
            for i, image in enumerate(images):
                page_num = start_page + i
                logger.info(f"Running OCR on page {page_num}...")
                
                page_text = await self._ocr_image(image)
                
                logger.info(f"Page {page_num} OCR result: {len(page_text)} characters")
                logger.info(f"Page {page_num} preview: {page_text[:200] if page_text else 'EMPTY'}...")
                
                batch_text += f"--- Page {page_num} ---\n{page_text}\n\n"
            
            self._log_operation("EXTRACT_PDF_PAGES_OCR", {
                "pages_processed": len(images),
                "total_text_length": len(batch_text),
                "text_preview": batch_text[:500] if batch_text else "EMPTY",
                "success": True
            }, is_request=False)
            
            return batch_text
            
        except Exception as e:
            logger.error(f"Failed to extract PDF pages {start_page}-{end_page}: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            
            self._log_operation("EXTRACT_PDF_PAGES_OCR", {
                "error": str(e),
                "success": False
            }, is_request=False)
            
            return ""

    async def extract_text_ai(self, file_path: str) -> str:
        """Extract text using AI-based method (direct image analysis)"""
        self._log_operation("EXTRACT_TEXT_AI", {"file_path": file_path})
        
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.pdf':
                text = await self._extract_pdf_with_ai_processing(file_path)
            elif file_ext in ['.jpg', '.jpeg', '.png']:
                text = await self._extract_image_with_ai(file_path)
            elif file_ext == '.pptx':
                text = await self._extract_pptx_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            self._log_operation("EXTRACT_TEXT_AI", {
                "text_length": len(text) if text else 0,
                "text_preview": (text[:500] + "...") if text and len(text) > 500 else text,
                "success": True
            }, is_request=False)
            
            return text
                
        except Exception as e:
            logger.error(f"Failed to extract text with AI from {file_path}: {str(e)}")
            self._log_operation("EXTRACT_TEXT_AI", {
                "error": str(e),
                "success": False
            }, is_request=False)
            return ""

    async def extract_text_ai_batched(self, file_path: str, session_id: str) -> List[TextBatch]:
        """
        Extract text using AI with batching support
        Returns list of TextBatch objects with extracted text
        """
        self._log_operation("EXTRACT_TEXT_AI_BATCHED", {
            "file_path": file_path,
            "session_id": session_id
        })
        
        try:
            # Create batches first - divides large files into manageable chunks
            batches = self.create_batches(file_path, session_id)
            logger.info(f"Created {len(batches)} batches for AI extraction")
            
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.pdf':
                # Convert PDF to images and process in batches - AI works better with images
                logger.info("Converting PDF to images for AI analysis...")
                images = pdf2image.convert_from_path(file_path, dpi=150)
                logger.info(f"Converted PDF to {len(images)} images")
                
                # Process each batch separately to manage memory usage
                for batch in batches:
                    start_page, end_page = batch.page_range
                    # Get images for this batch (0-indexed conversion)
                    batch_images = images[start_page-1:end_page]
                    
                    logger.info(f"Processing batch {batch.batch_number}: {len(batch_images)} images")
                    
                    # Save images temporarily for AI processing
                    temp_paths = []
                    for i, image in enumerate(batch_images):
                        temp_path = f"/tmp/batch_{batch.batch_id}_page_{start_page+i}.png"
                        image.save(temp_path)
                        temp_paths.append(temp_path)
                        logger.info(f"Saved temp image: {temp_path}")
                    
                    # Analyze with AI - send multiple images for context
                    batch_text = await self.ai_service.analyze_images(temp_paths)
                    batch.text_content = batch_text
                    
                    logger.info(f"AI analysis for batch {batch.batch_number}: {len(batch_text)} characters")
                    
                    # Clean up temp files to prevent disk space issues
                    for temp_path in temp_paths:
                        if os.path.exists(temp_path):
                            os.remove(temp_path)
                            
            elif file_ext in ['.jpg', '.jpeg', '.png']:
                # Single image - single batch processing
                text = await self._extract_image_with_ai(file_path)
                batches[0].text_content = text
            elif file_ext in ['.pptx', '.ppt']:
                # Extract text from presentation slides
                text = await self._extract_pptx_text(file_path)
                batches[0].text_content = text
            
            # Log processing statistics for monitoring
            total_text = sum(len(b.text_content) for b in batches)
            self._log_operation("EXTRACT_TEXT_AI_BATCHED", {
                "batches_processed": len(batches),
                "total_text_length": total_text,
                "success": True
            }, is_request=False)
            
            return batches
            
        except Exception as e:
            logger.error(f"Failed to extract text with AI batching: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            
            # Log error for debugging
            self._log_operation("EXTRACT_TEXT_AI_BATCHED", {
                "error": str(e),
                "success": False
            }, is_request=False)
            
            return []

    async def _extract_pdf_with_ocr_processing(self, file_path: str, session_id: str = None) -> str:
        """Enhanced PDF processing with page-based OCR"""
        self._log_operation("EXTRACT_PDF_WITH_OCR", {"file_path": file_path})
        
        try:
            # Get page count
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                page_count = len(pdf_reader.pages)
            
            logger.info(f"Processing PDF with {page_count} pages using OCR")
            
            # First try direct text extraction
            direct_text = await self._extract_pdf_text(file_path)
            logger.info(f"Direct PDF text extraction: {len(direct_text)} characters")
            
            # If direct extraction got sufficient text, use it
            if len(direct_text.strip()) > 100:
                logger.info("Using direct text extraction (sufficient text found)")
                self._log_operation("EXTRACT_PDF_WITH_OCR", {
                    "method": "direct",
                    "text_length": len(direct_text),
                    "success": True
                }, is_request=False)
                return direct_text
            
            # Otherwise, use OCR
            logger.info("Direct extraction insufficient, using OCR...")
            
            if page_count <= 5:
                # For PDFs with 5 or fewer pages, process normally
                text = await self._process_small_pdf(file_path)
            else:
                # For PDFs with more than 5 pages, use batched processing
                text = await self._process_large_pdf(file_path, page_count)
            
            self._log_operation("EXTRACT_PDF_WITH_OCR", {
                "method": "ocr",
                "page_count": page_count,
                "text_length": len(text) if text else 0,
                "text_preview": (text[:500] + "...") if text and len(text) > 500 else text,
                "success": True
            }, is_request=False)
            
            return text
                
        except Exception as e:
            logger.error(f"Failed to process PDF with OCR: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            
            self._log_operation("EXTRACT_PDF_WITH_OCR", {
                "error": str(e),
                "success": False
            }, is_request=False)
            
            return ""

    async def _extract_pdf_with_ai_processing(self, file_path: str) -> str:
        """Enhanced PDF processing with AI-based analysis"""
        self._log_operation("EXTRACT_PDF_WITH_AI", {"file_path": file_path})
        
        try:
            # Convert PDF pages to images
            logger.info("Converting PDF to images for AI processing...")
            images = pdf2image.convert_from_path(file_path, dpi=150)
            logger.info(f"Converted to {len(images)} images")
            
            if len(images) <= 5:
                # Process all images with AI
                text = await self._process_images_with_ai(images, batch_size=len(images))
            else:
                # Process in batches of 3 pages
                text = await self._process_images_with_ai(images, batch_size=3)
            
            self._log_operation("EXTRACT_PDF_WITH_AI", {
                "page_count": len(images),
                "text_length": len(text) if text else 0,
                "success": True
            }, is_request=False)
            
            return text
                
        except Exception as e:
            logger.error(f"Failed to process PDF with AI: {str(e)}")
            self._log_operation("EXTRACT_PDF_WITH_AI", {
                "error": str(e),
                "success": False
            }, is_request=False)
            return ""

    async def _process_small_pdf(self, file_path: str) -> str:
        """Process PDF with 5 or fewer pages"""
        logger.info("Processing small PDF (≤5 pages)")
        
        try:
            # Convert to images and OCR
            logger.info("Converting PDF to images...")
            images = pdf2image.convert_from_path(file_path, dpi=200)
            logger.info(f"Got {len(images)} images")
            
            ocr_text = ""
            for i, image in enumerate(images):
                logger.info(f"Running OCR on page {i+1}...")
                page_text = await self._ocr_image(image)
                logger.info(f"Page {i+1} OCR: {len(page_text)} chars - Preview: {page_text[:100] if page_text else 'EMPTY'}...")
                ocr_text += f"--- Page {i+1} ---\n{page_text}\n\n"
            
            logger.info(f"Total OCR text: {len(ocr_text)} characters")
            return ocr_text
            
        except Exception as e:
            logger.error(f"Failed to process small PDF: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            return ""

    async def _process_large_pdf(self, file_path: str, page_count: int) -> str:
        """Process PDF with more than 5 pages using batched OCR"""
        logger.info(f"Processing large PDF ({page_count} pages)")
        
        try:
            # Convert all pages to images
            logger.info("Converting PDF to images...")
            images = pdf2image.convert_from_path(file_path, dpi=200)
            logger.info(f"Got {len(images)} images")
            
            # Process in batches of 2 pages
            batch_size = 2
            all_text = ""
            
            for i in range(0, len(images), batch_size):
                batch_images = images[i:i + batch_size]
                batch_text = ""
                
                logger.info(f"Processing batch {i//batch_size + 1}: pages {i+1}-{min(i+batch_size, len(images))}")
                
                # OCR each image in the batch
                for j, image in enumerate(batch_images):
                    page_num = i + j + 1
                    logger.info(f"Running OCR on page {page_num}...")
                    page_text = await self._ocr_image(image)
                    logger.info(f"Page {page_num} OCR: {len(page_text)} chars")
                    batch_text += f"--- Page {page_num} ---\n{page_text}\n\n"
                
                all_text += batch_text
            
            logger.info(f"Total OCR text from large PDF: {len(all_text)} characters")
            return all_text
            
        except Exception as e:
            logger.error(f"Failed to process large PDF: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            return ""

    async def _process_images_with_ai(self, images: List[Image.Image], batch_size: int) -> str:
        """Process images with AI in batches"""
        logger.info(f"Processing {len(images)} images with AI in batches of {batch_size}")
        
        try:
            all_text = ""
            
            for i in range(0, len(images), batch_size):
                batch_images = images[i:i + batch_size]
                logger.info(f"AI processing batch: images {i+1}-{i+len(batch_images)}")
                
                # Convert images to base64 for AI processing
                image_paths = []
                for j, image in enumerate(batch_images):
                    # Save image temporarily
                    temp_path = f"/tmp/ai_page_{i+j+1}_{uuid.uuid4().hex[:8]}.png"
                    image.save(temp_path)
                    image_paths.append(temp_path)
                
                # Send to AI for analysis
                batch_text = await self.ai_service.analyze_images(image_paths)
                all_text += batch_text + "\n\n"
                
                logger.info(f"AI batch result: {len(batch_text)} characters")
                
                # Clean up temp files
                for temp_path in image_paths:
                    if os.path.exists(temp_path):
                        os.remove(temp_path)
            
            return all_text
            
        except Exception as e:
            logger.error(f"Failed to process images with AI: {str(e)}")
            return ""

    async def _ocr_image(self, image: Image.Image) -> str:
        """Perform OCR on a single image"""
        self._log_operation("OCR_IMAGE", {"image_size": image.size if image else None})
        
        try:
            # Convert PIL image to numpy array
            img_array = np.array(image)
            logger.info(f"Image array shape: {img_array.shape}")
            
            # Preprocess image
            processed_img = self._preprocess_image_array(img_array)
            logger.info(f"Preprocessed image shape: {processed_img.shape}")
            
            # Check if tesseract is available
            try:
                tesseract_version = pytesseract.get_tesseract_version()
                logger.info(f"Tesseract version: {tesseract_version}")
            except Exception as e:
                logger.error(f"Tesseract not available: {e}")
                raise Exception(f"Tesseract OCR not available: {e}")
            
            # Perform OCR with multiple configurations for better results
            logger.info("Running pytesseract OCR...")
            
            # Try with PSM 6 (uniform block of text)
            text = pytesseract.image_to_string(processed_img, config='--psm 6 --oem 3')
            
            # If result is too short, try PSM 3 (fully automatic)
            if len(text.strip()) < 50:
                logger.info("PSM 6 result too short, trying PSM 3...")
                text_psm3 = pytesseract.image_to_string(processed_img, config='--psm 3 --oem 3')
                if len(text_psm3.strip()) > len(text.strip()):
                    text = text_psm3
            
            # If still too short, try on original image without preprocessing
            if len(text.strip()) < 50:
                logger.info("Preprocessed OCR result too short, trying original image...")
                text_original = pytesseract.image_to_string(img_array, config='--psm 3 --oem 3')
                if len(text_original.strip()) > len(text.strip()):
                    text = text_original
            
            logger.info(f"OCR result: {len(text)} characters")
            logger.info(f"OCR preview: {text[:200] if text else 'EMPTY'}...")
            
            self._log_operation("OCR_IMAGE", {
                "text_length": len(text) if text else 0,
                "text_preview": (text[:300] + "...") if text and len(text) > 300 else text,
                "success": True
            }, is_request=False)
            
            return text.strip()
            
        except Exception as e:
            logger.error(f"Failed to OCR image: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            
            self._log_operation("OCR_IMAGE", {
                "error": str(e),
                "success": False
            }, is_request=False)
            
            return ""
    
    async def _extract_image_with_ai(self, file_path: str) -> str:
        """Extract text from image using AI"""
        self._log_operation("EXTRACT_IMAGE_WITH_AI", {"file_path": file_path})
        
        try:
            text = await self.ai_service.analyze_images([file_path])
            
            self._log_operation("EXTRACT_IMAGE_WITH_AI", {
                "text_length": len(text) if text else 0,
                "success": True
            }, is_request=False)
            
            return text
        except Exception as e:
            logger.error(f"Failed to extract image with AI: {str(e)}")
            self._log_operation("EXTRACT_IMAGE_WITH_AI", {
                "error": str(e),
                "success": False
            }, is_request=False)
            return ""

    async def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text directly from PDF"""
        self._log_operation("EXTRACT_PDF_TEXT_DIRECT", {"file_path": file_path})
        
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                logger.info(f"PDF has {len(pdf_reader.pages)} pages")
                
                for i, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    logger.info(f"Page {i+1} direct extraction: {len(page_text) if page_text else 0} chars")
                    text += page_text + "\n"
            
            self._log_operation("EXTRACT_PDF_TEXT_DIRECT", {
                "text_length": len(text) if text else 0,
                "text_preview": (text[:300] + "...") if text and len(text) > 300 else text,
                "success": True
            }, is_request=False)
            
            return text.strip()
        except Exception as e:
            logger.error(f"Failed to extract PDF text: {str(e)}")
            self._log_operation("EXTRACT_PDF_TEXT_DIRECT", {
                "error": str(e),
                "success": False
            }, is_request=False)
            return ""

    async def _extract_image_ocr(self, file_path: str) -> str:
        """Extract text from image using OCR"""
        self._log_operation("EXTRACT_IMAGE_OCR", {"file_path": file_path})
        
        try:
            # Load image
            image = Image.open(file_path)
            logger.info(f"Loaded image: {image.size}, mode: {image.mode}")
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # OCR the image
            text = await self._ocr_image(image)
            
            self._log_operation("EXTRACT_IMAGE_OCR", {
                "text_length": len(text) if text else 0,
                "success": True
            }, is_request=False)
            
            return text
            
        except Exception as e:
            logger.error(f"Failed to extract image with OCR: {str(e)}")
            self._log_operation("EXTRACT_IMAGE_OCR", {
                "error": str(e),
                "success": False
            }, is_request=False)
            return ""

    async def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PowerPoint files"""
        self._log_operation("EXTRACT_PPTX_TEXT", {"file_path": file_path})
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                text += slide_text + "\n"
            
            self._log_operation("EXTRACT_PPTX_TEXT", {
                "slides_count": len(prs.slides),
                "text_length": len(text),
                "success": True
            }, is_request=False)
            
            return text.strip()
            
        except Exception as e:
            logger.error(f"Failed to extract PPTX text: {str(e)}")
            self._log_operation("EXTRACT_PPTX_TEXT", {
                "error": str(e),
                "success": False
            }, is_request=False)
            return ""

    def _preprocess_image(self, image_path: str) -> np.ndarray:
        """Preprocess image for better OCR results"""
        try:
            # Load image
            image = cv2.imread(image_path)
            return self._preprocess_image_array(image)
            
        except Exception as e:
            logger.error(f"Failed to preprocess image: {str(e)}")
            return None

    def _preprocess_image_array(self, img_array: np.ndarray) -> np.ndarray:
        """Preprocess image array for better OCR results"""
        try:
            logger.info(f"Preprocessing image array with shape: {img_array.shape}")
            
            # Convert to grayscale if needed
            if len(img_array.shape) == 3:
                gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array
            
            logger.info(f"Grayscale shape: {gray.shape}")
            
            # Apply denoising
            denoised = cv2.fastNlMeansDenoising(gray, None, 10, 7, 21)
            
            # Apply adaptive thresholding for better results on varied backgrounds
            thresh = cv2.adaptiveThreshold(
                denoised, 255, 
                cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                cv2.THRESH_BINARY, 11, 2
            )
            
            logger.info(f"Preprocessed image shape: {thresh.shape}")
            
            return thresh
            
        except Exception as e:
            logger.error(f"Failed to preprocess image array: {str(e)}")
            return img_array
